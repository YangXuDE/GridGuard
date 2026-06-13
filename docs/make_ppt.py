"""Generate GridGuard_Presentation.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

# ── palette ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black
BG2      = RGBColor(0x13, 0x1A, 0x24)   # card bg
BLUE     = RGBColor(0x38, 0xBD, 0xF8)   # primary accent
GREEN    = RGBColor(0x34, 0xD3, 0x99)   # success / secure
AMBER    = RGBColor(0xFB, 0xBF, 0x24)   # warning
RED      = RGBColor(0xF8, 0x71, 0x71)   # danger
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
BORDER   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT   = RGBColor(0xA7, 0x8B, 0xFA)   # purple accent

W, H = Inches(13.33), Inches(7.5)       # 16:9

SCREENSHOTS = {
    "idle":    "screenshots/dashboard-idle.png",
    "result":  "screenshots/dashboard-result.png",
    "actions": "screenshots/dashboard-actions.png",
}

# ── helpers ────────────────────────────────────────────────────────────────────
def solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill=None, border=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = Pt(0)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    return shape

def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=Pt(16), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def pill(slide, text, l, t, color, text_color=None, w=None):
    """Small rounded-rect badge."""
    from pptx.util import Inches, Pt
    if w is None:
        w = Inches(1.6)
    h = Inches(0.28)
    r = add_rect(slide, l, t, w, h, fill=color)
    r.line.width = Pt(0)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color or BG
    return r

def divider(slide, y):
    add_rect(slide, Inches(0.5), y, Inches(12.33), Inches(0.018), fill=BORDER)

# ── slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    solid_bg(slide, BG)

    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=BLUE)

    # glow circle decoration
    add_rect(slide, Inches(9.5), Inches(-1), Inches(5), Inches(5), fill=RGBColor(0x0F,0x23,0x3A))

    # tag
    pill(slide, "E.ON HACKATHON  ·  ENERGY × AI", Inches(0.5), Inches(1.5),
         BLUE, WHITE, w=Inches(4.2))

    # headline
    add_text(slide, "GridGuard", Inches(0.5), Inches(2.0), Inches(8), Inches(1.4),
             size=Pt(72), bold=True, color=WHITE)
    add_text(slide, "AI-Powered Grid Security\n& Contingency Management",
             Inches(0.5), Inches(3.3), Inches(8), Inches(1.2),
             size=Pt(26), bold=False, color=MUTED)

    # status pill
    pill(slide, "✦  N-1 Security Restored", Inches(0.5), Inches(4.8),
         GREEN, BG, w=Inches(2.8))

    # team
    add_text(slide, "Mengyu Zhang  ·  Chen Zhao  ·  Cici  ·  Yang Xu  ·  Weiting Liang",
             Inches(0.5), Inches(6.6), Inches(10), Inches(0.5),
             size=Pt(13), color=MUTED)

    # screenshot on right
    img = os.path.join(os.path.dirname(__file__), SCREENSHOTS["idle"])
    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(7.6), Inches(1.2),
                                  Inches(5.4), Inches(3.0))


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=RED)

    add_text(slide, "The Problem", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    boxes = [
        (RED,   "⚡  N-1 Contingency Risk",
                "A single line trip can cascade into\nwider outages if not corrected\nwithin minutes."),
        (AMBER, "🕒  Operator Overload",
                "Modern grids with high DER penetration\nproduce thousands of contingency\ncombinations — too many to screen manually."),
        (BLUE,  "📊  No Physics-Aware AI",
                "Existing tools are either pure\noptimisation (no explainability) or\npure heuristics (no physics)."),
    ]
    for i, (col, title, body) in enumerate(boxes):
        lft = Inches(0.5 + i * 4.22)
        add_rect(slide, lft, Inches(1.3), Inches(3.9), Inches(5.5),
                 fill=BG2, border=col, border_w=Pt(1.5))
        pill(slide, title, lft + Inches(0.18), Inches(1.55), col, BG,
             w=Inches(3.55))
        add_text(slide, body, lft + Inches(0.18), Inches(2.05),
                 Inches(3.55), Inches(4.6),
                 size=Pt(15), color=MUTED, wrap=True)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=GREEN)

    add_text(slide, "Our Solution", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    add_text(slide,
             "GridGuard is a real-time control-room decision-support dashboard that combines "
             "physics-based N-1 screening with an LLM reasoning loop to restore grid security "
             "at the lowest corrective cost.",
             Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.9),
             size=Pt(17), color=MUTED, wrap=True)

    steps = [
        (BLUE,   "1",  "Two-Stage N-1 Screen",
                       "DC power-flow pre-screen all 123 lines → AC solve risky shortlist. "
                       "Risk score = overload severity × bus count × islanding penalty."),
        (ACCENT, "2",  "Corrective Action Space",
                       "Build feasible DER options (BESS / RES curtail / load shed) with exact "
                       "sensitivity (Δloading pp/MW) and economic cost (€/MWh)."),
        (GREEN,  "3",  "DeepSeek LLM Agent",
                       "Context-aware reasoning loop: batteries first → curtail RES → shed load. "
                       "Iterates until N-1 secure. Benchmarked against LP (SCOPF-lite) optimum."),
        (AMBER,  "4",  "Operator Approval",
                       "Full explainability — every dispatch decision shown with cost & sensitivity. "
                       "AI recommends; operator accepts or overrides."),
    ]
    for i, (col, num, title, body) in enumerate(steps):
        row, col_idx = divmod(i, 2)
        lft = Inches(0.5 + col_idx * 6.42)
        top = Inches(2.35 + row * 2.35)
        add_rect(slide, lft, top, Inches(6.1), Inches(2.1), fill=BG2, border=col, border_w=Pt(1.2))
        # number circle
        circ = add_rect(slide, lft + Inches(0.18), top + Inches(0.18),
                        Inches(0.42), Inches(0.42), fill=col)
        add_text(slide, num, lft + Inches(0.18), top + Inches(0.16),
                 Inches(0.42), Inches(0.42),
                 size=Pt(14), bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(slide, title, lft + Inches(0.72), top + Inches(0.18),
                 Inches(5.2), Inches(0.38),
                 size=Pt(15), bold=True, color=WHITE)
        add_text(slide, body, lft + Inches(0.18), top + Inches(0.65),
                 Inches(5.75), Inches(1.3),
                 size=Pt(13), color=MUTED, wrap=True)


def slide_screenshot(prs, key, headline, sub):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, RGBColor(0x08, 0x0D, 0x14))
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=BLUE)

    add_text(slide, headline, Inches(0.3), Inches(0.18), Inches(12.5), Inches(0.55),
             size=Pt(26), bold=True, color=WHITE)
    add_text(slide, sub, Inches(0.3), Inches(0.72), Inches(12.5), Inches(0.4),
             size=Pt(14), color=MUTED)
    divider(slide, Inches(1.12))

    img = os.path.join(os.path.dirname(__file__), SCREENSHOTS[key])
    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(0.25), Inches(1.2),
                                  Inches(12.83), Inches(6.1))


def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=ACCENT)

    add_text(slide, "Technology Stack", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    left_items = [
        (BLUE,   "React 19 + TypeScript",    "Type-safe component tree"),
        (BLUE,   "TanStack Router / Start",  "File-based routing, SSR-ready"),
        (BLUE,   "Vite 8",                   "Sub-second HMR, optimised build"),
        (ACCENT, "Tailwind CSS v4",          "Design tokens, dark-mode native"),
        (ACCENT, "Radix UI",                 "Accessible headless primitives"),
        (GREEN,  "Recharts",                 "Bar/line charts for cost analysis"),
    ]
    right_items = [
        (AMBER,  "pandapower",               "IEEE 123-node N-1 physics engine"),
        (AMBER,  "Two-stage screening",      "DC pre-screen → AC solve shortlist"),
        (RED,    "DeepSeek LLM",             "Context-aware corrective loop"),
        (GREEN,  "LP Baseline (SCOPF-lite)", "Mathematical cost optimum reference"),
        (MUTED,  "Static TypeScript data",   "Pre-simulated scenarios, zero backend"),
    ]

    for i, (col, title, sub) in enumerate(left_items):
        top = Inches(1.25 + i * 0.96)
        pill(slide, title, Inches(0.5), top + Inches(0.05), col, BG, w=Inches(2.8))
        add_text(slide, sub, Inches(3.45), top + Inches(0.02),
                 Inches(2.7), Inches(0.45), size=Pt(13), color=MUTED)

    for i, (col, title, sub) in enumerate(right_items):
        top = Inches(1.25 + i * 0.96)
        pill(slide, title, Inches(6.7), top + Inches(0.05), col, BG, w=Inches(3.0))
        add_text(slide, sub, Inches(9.85), top + Inches(0.02),
                 Inches(3.1), Inches(0.45), size=Pt(13), color=MUTED)

    # vertical divider
    add_rect(slide, Inches(6.42), Inches(1.2), Inches(0.018), Inches(5.8), fill=BORDER)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=GREEN)

    add_text(slide, "Key Results", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    kpis = [
        (GREEN,  "119",  "Contingencies\nScreened",    "All 123-node lines evaluated in < 1 s"),
        (RED,    "3",    "Catastrophic\nContingencies", "Immediately escalated to operator"),
        (BLUE,   "2",    "Corrective\nIterations",      "LLM restores security in ≤ 2 steps"),
        (AMBER,  "96.2%","Final Max\nLine Loading",     "Down from 136% post-fault"),
        (ACCENT, "€1 211","Corrective\nCost /h",        "vs €980 LP optimum — 24% gap"),
    ]
    for i, (col, val, label, note) in enumerate(kpis):
        lft = Inches(0.35 + i * 2.55)
        add_rect(slide, lft, Inches(1.3), Inches(2.35), Inches(3.8),
                 fill=BG2, border=col, border_w=Pt(1.5))
        add_text(slide, val, lft + Inches(0.15), Inches(1.55),
                 Inches(2.05), Inches(0.9),
                 size=Pt(38), bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, label, lft + Inches(0.15), Inches(2.45),
                 Inches(2.05), Inches(0.65),
                 size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, note, lft + Inches(0.1), Inches(3.15),
                 Inches(2.15), Inches(0.85),
                 size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

    add_text(slide,
             "Distribution-feeder dispatch priority:  BESS Discharge  →  Curtail Renewables  →  Shed Load",
             Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
             size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide,
             "(cheapest to most expensive per MW of overload relief)",
             Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.4),
             size=Pt(12), color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)


def slide_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=AMBER)

    add_text(slide, "Three Demo Scenarios", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    scenarios = [
        (RED,   "Heatwave Peak",  "N-1: Line 13 trip",
                "130 % load  ·  16:00  ·  34 °C",
                "12 overloads  ·  Risk Score 163.7  →  Catastrophic",
                "BESS @ Bus 51 discharge 0.39 MW  +  load curtail"),
        (BLUE,  "Solar Midday",   "N-1: Line 19 trip",
                "100 % load  ·  12:00  ·  clear sky",
                "8 overloads  ·  Risk Score 102.5  →  Dangerous",
                "Curtail wind farm @ Bus 111  +  BESS @ Bus 76"),
        (AMBER, "Storm Evening",  "N-1: Line 55 trip",
                "120 % load  ·  18:00  ·  storm",
                "5 overloads  ·  Risk Score 88.4  →  Elevated",
                "BESS discharge  +  load curtail Bus 47"),
    ]
    for i, (col, name, fault, cond, screen, action) in enumerate(scenarios):
        lft = Inches(0.5 + i * 4.22)
        add_rect(slide, lft, Inches(1.3), Inches(3.9), Inches(5.8),
                 fill=BG2, border=col, border_w=Pt(1.5))
        pill(slide, name, lft + Inches(0.18), Inches(1.5), col, BG, w=Inches(3.55))
        add_text(slide, fault, lft + Inches(0.18), Inches(1.95),
                 Inches(3.55), Inches(0.4), size=Pt(14), bold=True, color=WHITE)
        add_text(slide, cond, lft + Inches(0.18), Inches(2.38),
                 Inches(3.55), Inches(0.38), size=Pt(12), color=MUTED)
        divider_mini = add_rect(slide, lft + Inches(0.18), Inches(2.82),
                                Inches(3.55), Inches(0.018), fill=BORDER)
        add_text(slide, "Screening result:", lft + Inches(0.18), Inches(2.9),
                 Inches(3.55), Inches(0.32), size=Pt(11), bold=True, color=MUTED)
        add_text(slide, screen, lft + Inches(0.18), Inches(3.22),
                 Inches(3.55), Inches(0.55), size=Pt(12), color=col, wrap=True)
        add_text(slide, "AI corrective action:", lft + Inches(0.18), Inches(3.85),
                 Inches(3.55), Inches(0.32), size=Pt(11), bold=True, color=MUTED)
        add_text(slide, action, lft + Inches(0.18), Inches(4.17),
                 Inches(3.55), Inches(0.75), size=Pt(12), color=WHITE, wrap=True)


def slide_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=ACCENT)

    add_text(slide, "The Team", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             size=Pt(36), bold=True, color=WHITE)
    divider(slide, Inches(1.05))

    members = [
        "Mengyu Zhang",
        "Chen Zhao",
        "Cici",
        "Yang Xu",
        "Weiting Liang",
    ]
    cols = [BLUE, GREEN, ACCENT, AMBER, RED]
    for i, (name, col) in enumerate(zip(members, cols)):
        lft = Inches(0.45 + i * 2.5)
        add_rect(slide, lft, Inches(1.6), Inches(2.2), Inches(2.6),
                 fill=BG2, border=col, border_w=Pt(1.5))
        # avatar circle
        circ = add_rect(slide, lft + Inches(0.72), Inches(1.85),
                        Inches(0.76), Inches(0.76), fill=col)
        initials = "".join(w[0].upper() for w in name.split()[:2])
        add_text(slide, initials,
                 lft + Inches(0.72), Inches(1.85),
                 Inches(0.76), Inches(0.76),
                 size=Pt(20), bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(slide, name, lft + Inches(0.1), Inches(2.72),
                 Inches(2.0), Inches(0.55),
                 size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Built for the E.ON Hackathon 2025  ·  Energy × AI Grid Operation Agents Track",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), fill=BLUE)
    add_rect(slide, Inches(0), H - Inches(0.12), W, Inches(0.12), fill=BLUE)

    add_text(slide, "GridGuard", Inches(1.5), Inches(2.0), Inches(10), Inches(1.4),
             size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Physics-informed · Explainable · Operator-approved",
             Inches(1.5), Inches(3.35), Inches(10), Inches(0.7),
             size=Pt(22), color=MUTED, align=PP_ALIGN.CENTER)

    pill(slide, "✦  N-1 Security Restored  ✦",
         Inches(4.5), Inches(4.3), GREEN, BG, w=Inches(4.33))

    add_text(slide, "github.com/YangXuDE/GridGuard",
             Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)


# ── main ───────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_screenshot(prs, "result",
                     "N-1 Violation Detected — Risk Screening",
                     "Line L13 trip · 12 simultaneous overloads · Risk Score 163.7 (Catastrophic)")
    slide_screenshot(prs, "actions",
                     "Corrective Action Space & Operator Action Plan",
                     "AI recommends BESS @ Bus 76 discharge — best sensitivity-to-cost ratio")
    slide_scenarios(prs)
    slide_tech(prs)
    slide_results(prs)
    slide_team(prs)
    slide_closing(prs)

    out = os.path.join(os.path.dirname(__file__), "GridGuard_Presentation.pptx")
    prs.save(out)
    print(f"Saved → {out}")

if __name__ == "__main__":
    main()
